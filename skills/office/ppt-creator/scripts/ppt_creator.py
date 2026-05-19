#!/usr/bin/env python3
"""
PowerPoint 演示文稿生成器

用法:
    python ppt_creator.py --title "季度报告" --output report.pptx
    python ppt_creator.py --slides title,content,content --title "季度报告" --output report.pptx
    python ppt_creator.py --dry-run --title "测试" --slides title,content
"""

import sys
import json
import argparse
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

# 默认配色
THEME_COLORS = {
    "primary": RGBColor(0x2C, 0x5F, 0x8A),
    "secondary": RGBColor(0x4A, 0x90, 0xD9),
    "accent": RGBColor(0xF5, 0xA6, 0x23),
    "text": RGBColor(0x33, 0x33, 0x33),
    "light_text": RGBColor(0x66, 0x66, 0x66),
    "white": RGBColor(0xFF, 0xFF, 0xFF),
    "bg_light": RGBColor(0xF5, 0xF5, 0xF5),
}


def create_presentation(title="", slides=None, output_path=None, dry_run=False):
    """创建 PowerPoint 演示文稿"""
    if not HAS_PPTX:
        return {"error": "python-pptx 未安装，请运行: pip install python-pptx"}

    if dry_run:
        return {"status": "dry_run", "title": title, "slides": slides or [], "message": "仅模拟运行，未生成文件"}

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide_configs = slides or ["title", "content"]
    created_slides = []

    for i, slide_type in enumerate(slide_configs):
        if slide_type == "title":
            slide_layout = prs.slide_layouts[6]  # Blank
            slide = prs.slides.add_slide(slide_layout)
            add_title_slide(slide, title or "演示文稿")
            created_slides.append({"index": i, "type": "title", "title": title})
        elif slide_type == "content":
            slide_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(slide_layout)
            add_content_slide(slide, f"第 {i} 页")
            created_slides.append({"index": i, "type": "content"})
        elif slide_type == "section":
            slide_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(slide_layout)
            add_section_slide(slide, f"章节 {i}")
            created_slides.append({"index": i, "type": "section"})
        elif slide_type == "two_column":
            slide_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(slide_layout)
            add_two_column_slide(slide, f"第 {i} 页")
            created_slides.append({"index": i, "type": "two_column"})

    if output_path:
        prs.save(output_path)
        return {
            "status": "success",
            "output": str(Path(output_path).resolve()),
            "slides": len(created_slides),
            "slide_details": created_slides,
        }
    else:
        return {"error": "请指定 --output 参数"}


def add_title_slide(slide, title_text):
    """添加标题页"""
    # 背景色块
    from pptx.util import Emu
    shapes = slide.shapes
    bg_shape = shapes.add_shape(
        1, Emu(0), Emu(0), Emu(12192000), Emu(6858000)  # 13.333in x 7.5in
    )
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = THEME_COLORS["primary"]
    bg_shape.line.fill.background()

    # 标题
    txBox = shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.333), Inches(2))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(44)
    p.font.color.rgb = THEME_COLORS["white"]
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    # 副标题
    sub_box = shapes.add_textbox(Inches(1), Inches(4.8), Inches(11.333), Inches(1))
    sub_tf = sub_box.text_frame
    sub_p = sub_tf.paragraphs[0]
    sub_p.text = "自动生成"
    sub_p.font.size = Pt(24)
    sub_p.font.color.rgb = THEME_COLORS["white"]
    sub_p.alignment = PP_ALIGN.CENTER


def add_content_slide(slide, title_text):
    """添加内容页"""
    shapes = slide.shapes

    # 顶部色条
    header = shapes.add_shape(1, Emu(0), Emu(0), Emu(12192000), Emu(914400))
    header.fill.solid()
    header.fill.fore_color.rgb = THEME_COLORS["primary"]
    header.line.fill.background()

    # 标题
    txBox = shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(0.6))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(32)
    p.font.color.rgb = THEME_COLORS["white"]
    p.font.bold = True

    # 内容区域
    content_box = shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(11.5), Inches(5.5))
    content_tf = content_box.text_frame
    content_tf.word_wrap = True
    content_p = content_tf.paragraphs[0]
    content_p.text = "在此添加内容..."
    content_p.font.size = Pt(18)
    content_p.font.color.rgb = THEME_COLORS["text"]


def add_section_slide(slide, section_title):
    """添加章节页"""
    shapes = slide.shapes

    # 背景
    bg = shapes.add_shape(1, Emu(0), Emu(0), Emu(12192000), Emu(6858000))
    bg.fill.solid()
    bg.fill.fore_color.rgb = THEME_COLORS["bg_light"]
    bg.line.fill.background()

    # 左侧色条
    bar = shapes.add_shape(1, Emu(0), Emu(0), Emu(1828800), Emu(6858000))
    bar.fill.solid()
    bar.fill.fore_color.rgb = THEME_COLORS["primary"]
    bar.line.fill.background()

    # 标题
    txBox = shapes.add_textbox(Inches(2.5), Inches(3), Inches(8), Inches(1.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = section_title
    p.font.size = Pt(40)
    p.font.color.rgb = THEME_COLORS["primary"]
    p.font.bold = True


def add_two_column_slide(slide, title_text):
    """添加双栏内容页"""
    shapes = slide.shapes

    # 顶部色条
    header = shapes.add_shape(1, Emu(0), Emu(0), Emu(12192000), Emu(914400))
    header.fill.solid()
    header.fill.fore_color.rgb = THEME_COLORS["primary"]
    header.line.fill.background()

    # 标题
    txBox = shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(0.6))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(32)
    p.font.color.rgb = THEME_COLORS["white"]
    p.font.bold = True

    # 左栏
    left_box = shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(5.8), Inches(5.5))
    left_tf = left_box.text_frame
    left_tf.word_wrap = True
    left_p = left_tf.paragraphs[0]
    left_p.text = "左栏内容..."
    left_p.font.size = Pt(16)
    left_p.font.color.rgb = THEME_COLORS["text"]

    # 右栏
    right_box = shapes.add_textbox(Inches(7), Inches(1.3), Inches(5.8), Inches(5.5))
    right_tf = right_box.text_frame
    right_tf.word_wrap = True
    right_p = right_tf.paragraphs[0]
    right_p.text = "右栏内容..."
    right_p.font.size = Pt(16)
    right_p.font.color.rgb = THEME_COLORS["text"]


def main():
    parser = argparse.ArgumentParser(description="PowerPoint 演示文稿生成器")
    parser.add_argument("--title", default="演示文稿", help="演示文稿标题")
    parser.add_argument("--output", "-o", required=False, help="输出文件路径 (.pptx)")
    parser.add_argument(
        "--slides",
        default="title,content,content",
        help="幻灯片类型列表 (title,content,section,two_column)，逗号分隔",
    )
    parser.add_argument("--dry-run", action="store_true", help="仅模拟运行")
    parser.add_argument("--json", action="store_true", help="输出 JSON 格式")

    args = parser.parse_args()
    slides = [s.strip() for s in args.slides.split(",")]

    result = create_presentation(
        title=args.title,
        slides=slides,
        output_path=args.output,
        dry_run=args.dry_run,
    )

    if args.json:
        print(json.dumps(result, ensure_ascii=False, indent=2))
    else:
        if result.get("error"):
            print(f"❌ {result['error']}", file=sys.stderr)
            sys.exit(1)
        elif result.get("status") == "dry_run":
            print(f"📊 模拟生成: {result['title']}")
            print(f"   幻灯片: {', '.join(result['slides'])}")
        else:
            print(f"✅ 已生成: {result['output']}")
            print(f"   共 {result['slides']} 页幻灯片")


if __name__ == "__main__":
    main()
